"""
PowerPoint — outline via the cloud brain (JSON), deck built with
python-pptx (title slide, bullet slides, closing), then opened visibly
in PowerPoint via win32com.
"""
import re
from pathlib import Path

from config import Config
from brain.prompts import PPT_OUTLINE_PROMPT


def _slug(text):
    s = re.sub(r"[^\w\s-]", "", text)[:60].strip().replace(" ", "_")
    return s or "presentation"


def _open_in_powerpoint(path, ctx):
    try:
        import win32com.client as win32
        ppt = win32.DispatchEx("PowerPoint.Application")
        ppt.Visible = True
        ppt.Presentations.Open(str(path), WithWindow=True)

        def _closer(p=ppt):
            try:
                p.Quit()
            except Exception:
                pass
        ctx.registry.register("document", Path(path).name,
                              window_title=Path(path).stem, closer=_closer,
                              extra={"path": str(path)})
        return True
    except Exception:
        try:
            import os
            os.startfile(str(path))
            ctx.registry.register("document", Path(path).name,
                                  window_title=Path(path).stem,
                                  extra={"path": str(path)})
        except Exception:
            pass
        return False


def create_presentation(topic, slides, ctx):
    topic = (topic or "").strip()
    if not topic:
        return "A presentation about what, sir?"
    if not ctx.llm.available:
        return "My design brain needs the OpenRouter key, sir."

    try:
        n = int(slides) if str(slides).strip() else 8
    except Exception:
        n = 8
    n = max(3, min(n, 20))

    ctx.speaker.speak(f"Outlining a {n}-slide deck on {topic}, sir.")
    data = ctx.llm.quick_json(
        PPT_OUTLINE_PROMPT.format(topic=topic, n=n), max_tokens=2500)
    if not data or not data.get("slides"):
        return "I couldn't outline that presentation, sir."

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]
        title_layout = prs.slide_layouts[0]
        content_layout = prs.slide_layouts[1]

        ACCENT = RGBColor(0x1F, 0x4E, 0x78)
        DARK = RGBColor(0x22, 0x22, 0x22)

        slides_data = data["slides"][:n]
        for idx, spec in enumerate(slides_data):
            title = str(spec.get("title", "")).strip()
            bullets = [str(b).strip() for b in (spec.get("bullets") or []) if str(b).strip()]

            if idx == 0:
                # title slide
                slide = prs.slides.add_slide(title_layout)
                slide.shapes.title.text = title or str(data.get("title", topic))
                try:
                    slide.placeholders[1].text = bullets[0] if bullets else ""
                except Exception:
                    pass
                continue

            if idx == len(slides_data) - 1 and not bullets:
                # closing slide
                slide = prs.slides.add_slide(blank)
                box = slide.shapes.add_textbox(Inches(1), Inches(2.8),
                                               Inches(11.3), Inches(1.6))
                tf = box.text_frame
                tf.text = title or "Thank you"
                tf.paragraphs[0].font.size = Pt(48)
                tf.paragraphs[0].font.bold = True
                tf.paragraphs[0].font.color.rgb = ACCENT
                continue

            slide = prs.slides.add_slide(content_layout)
            slide.shapes.title.text = title
            try:
                slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = ACCENT
            except Exception:
                pass
            body = slide.placeholders[1].text_frame
            body.clear()
            for i, b in enumerate(bullets[:6]):
                para = body.paragraphs[0] if i == 0 else body.add_paragraph()
                para.text = b
                para.font.size = Pt(20)
                para.font.color.rgb = DARK
                para.level = 0

        title_txt = str(data.get("title", topic))
        path = Config.DESKTOP_PATH / f"{_slug(title_txt)}.pptx"
        prs.save(str(path))
    except Exception as exc:
        return f"Deck construction failed: {exc}."

    _open_in_powerpoint(path, ctx)
    return (f"Done, sir. {path.name} is on your desktop and open in "
            f"PowerPoint — {len(data['slides'][:n])} slides.")


# ---------------------------------------------------------------------------
# Skill dispatch entry
# ---------------------------------------------------------------------------
def handle(intent, ctx):
    skill = intent.get("skill")
    params = intent.get("params", {}) or {}
    if skill == "ppt.create":
        return create_presentation(params.get("topic", ""),
                                   params.get("slides", ""), ctx)
    return None
